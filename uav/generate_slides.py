"""
OPENDRONEID_IMPLEMENTATION.md → PowerPoint スライド生成
  - サイズ  : 4:3  (10" × 7.5")
  - 背景    : 白
  - フォント: 游ゴシック (Yu Gothic)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── サイズ ──────────────────────────────────────────────────────────────────
W = Inches(10)
H = Inches(7.5)
M = Inches(0.45)          # 左右マージン
IW = W - M * 2            # 使用幅 = 9.1"

# ── カラーパレット（白ベース）────────────────────────────────────────────────
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)   # 背景: 白
C_NAVY      = RGBColor(0x1A, 0x23, 0x7E)   # 濃紺: 見出し
C_BLUE      = RGBColor(0x19, 0x76, 0xD2)   # 青: アクセント
C_TEAL      = RGBColor(0x00, 0x89, 0x7B)   # 青緑: サブアクセント
C_TEXT      = RGBColor(0x21, 0x21, 0x21)   # 本文
C_SUBTEXT   = RGBColor(0x61, 0x61, 0x61)   # 補足テキスト
C_RULE      = RGBColor(0xBD, 0xC8, 0xE0)   # 区切り線
C_CODE_BG   = RGBColor(0xF3, 0xF6, 0xFB)   # コード背景
C_CODE_TEXT = RGBColor(0x1A, 0x20, 0x2C)   # コード文字
C_TBLH_BG   = RGBColor(0x1A, 0x23, 0x7E)   # テーブルヘッダ背景
C_TBLH_TEXT = RGBColor(0xFF, 0xFF, 0xFF)   # テーブルヘッダ文字
C_TBL_ROW1  = RGBColor(0xFF, 0xFF, 0xFF)   # テーブル奇数行
C_TBL_ROW2  = RGBColor(0xEE, 0xF2, 0xFF)   # テーブル偶数行
C_TBL_TEXT  = RGBColor(0x21, 0x21, 0x21)   # テーブル文字

FONT = "Yu Gothic"


# ── 基本ヘルパー ────────────────────────────────────────────────────────────
def rect(slide, left, top, width, height, color, border=False):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.width = Pt(0.5)
        s.line.color.rgb = C_RULE
    else:
        s.line.fill.background()
    return s


def set_white_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG


def txb(slide, text, left, top, width, height,
        size=14, bold=False, color=C_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return box


def bullets(slide, lines, left, top, width, height, size=13, color=C_TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT
        p.space_after = Pt(4)


def code(slide, text, left, top, width, height, size=10):
    rect(slide, left, top, width, height, C_CODE_BG, border=True)
    rect(slide, left, top, Inches(0.055), height, C_BLUE)
    box = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.1),
        width - Inches(0.18), height - Inches(0.2)
    )
    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = C_CODE_TEXT
        r.font.name = "Courier New"


def table(slide, headers, rows, left, top, width, col_widths=None, font_size=11):
    nr = len(rows) + 1
    nc = len(headers)
    row_h = Inches(0.38)
    tbl = slide.shapes.add_table(nr, nc, left, top, width, row_h * nr).table
    tbl.first_row = True

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    else:
        for col in tbl.columns:
            col.width = width // nc

    def cell_style(cell, text, header=False, even=False):
        cell.fill.solid()
        if header:
            cell.fill.fore_color.rgb = C_TBLH_BG
        elif even:
            cell.fill.fore_color.rgb = C_TBL_ROW2
        else:
            cell.fill.fore_color.rgb = C_TBL_ROW1

        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(font_size)
        r.font.bold = header
        r.font.color.rgb = C_TBLH_TEXT if header else C_TBL_TEXT
        r.font.name = FONT

    for ci, h in enumerate(headers):
        cell_style(tbl.cell(0, ci), h, header=True)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell_style(tbl.cell(ri + 1, ci), val, even=(ri % 2 == 1))


# ── スライドレイアウト共通 ───────────────────────────────────────────────────
def header(slide, title, sub=None):
    set_white_bg(slide)
    # 上部カラーバー
    rect(slide, 0, 0, W, Inches(0.07), C_NAVY)
    # 左縦アクセント
    rect(slide, 0, Inches(0.07), Inches(0.07), H - Inches(0.07), C_BLUE)
    # タイトル
    txb(slide, title,
        M, Inches(0.12), IW, Inches(0.55),
        size=22, bold=True, color=C_NAVY)
    # サブタイトル
    if sub:
        txb(slide, sub,
            M, Inches(0.65), IW, Inches(0.32),
            size=11, color=C_SUBTEXT)
    # 区切り線
    rect(slide, M, Inches(0.98), IW, Pt(1.2), C_RULE)


# ページ番号（右下）
def page_num(slide, n):
    txb(slide, str(n),
        W - Inches(0.55), H - Inches(0.32), Inches(0.4), Inches(0.25),
        size=9, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)


# ── プレゼンテーション生成 ────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ── スライド 1: タイトル ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_white_bg(s)

    # 上部ネイビーバンド
    rect(s, 0, 0, W, Inches(1.4), C_NAVY)
    txb(s, "zkp-m1 project",
        M, Inches(0.2), IW, Inches(0.35),
        size=10, color=RGBColor(0xB0, 0xBE, 0xFF))

    # タイトルエリア（白地）
    txb(s, "OpenDroneID (ASTM F3411-22a)",
        M, Inches(1.7), IW, Inches(0.8),
        size=30, bold=True, color=C_NAVY)
    txb(s, "実装解説",
        M, Inches(2.45), IW, Inches(0.6),
        size=24, bold=True, color=C_BLUE)

    # 区切り線
    rect(s, M, Inches(3.15), Inches(1.5), Pt(2.5), C_TEAL)

    txb(s, "BLE ブロードキャスト・Flutter 受信・ZK 証明との統合",
        M, Inches(3.4), IW, Inches(0.4),
        size=13, color=C_SUBTEXT)

    # 下部カラーバー
    rect(s, 0, H - Inches(0.55), W, Inches(0.55), C_NAVY)
    txb(s, "OpenDroneID / ASTM F3411-22a / ZK-SNARKs / Flutter / Raspberry Pi",
        M, H - Inches(0.48), IW, Inches(0.38),
        size=9, color=RGBColor(0xB0, 0xBE, 0xFF))

    # ── スライド 2: プロジェクト概要 ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "プロジェクト概要",
           "zkp-m1 — ドローンの接近をゼロ知識証明で示すデモシステム")
    page_num(s, 2)

    txb(s, "何を証明するか",
        M, Inches(1.12), IW, Inches(0.32),
        size=12, bold=True, color=C_NAVY)
    rect(s, M, Inches(1.12), Inches(0.055), Inches(0.32), C_BLUE)
    bullets(s, [
        "「ドローンが利用者の 30m 以内に、かつ 5 秒以内の時刻差で接近した」という事実を",
        "生の GPS 座標・飛行ルートを一切公開せずに、暗号学的に証明する",
    ], M + Inches(0.12), Inches(1.12), IW - Inches(0.15), Inches(0.65), size=12)

    # 3列ボックス（ドローン / 証明 / 利用者）
    for i, (title_t, color_t, lines) in enumerate([
        ("ドローン側", C_NAVY, [
            "Raspberry Pi",
            "OpenDroneID BLE 送信",
            "(ASTM F3411-22a)",
            "lat / lon / alt / time",
            "1秒ごとにブロードキャスト",
        ]),
        ("ZK 証明エンジン", C_BLUE, [
            "Circom 回路 (Groth16)",
            "Witness 計算 (.wcd)",
            "Proof 生成 (.zkey)",
            "Verify (verification_key)",
            "座標は private input",
        ]),
        ("利用者側", C_TEAL, [
            "Android (Flutter)",
            "BLE 受信 → ドローン位置",
            "GPS → 利用者位置",
            "カメラ録画 → 映像証拠",
            "端末上でオンデバイス証明",
        ]),
    ]):
        bx = M + Inches(3.05 * i)
        rect(s, bx, Inches(2.05), Inches(2.85), Inches(0.38), color_t)
        txb(s, title_t, bx + Inches(0.08), Inches(2.1), Inches(2.7), Inches(0.3),
            size=12, bold=True, color=C_WHITE)
        rect(s, bx, Inches(2.43), Inches(2.85), Inches(2.85), C_CODE_BG, border=True)
        bullets(s, lines, bx + Inches(0.12), Inches(2.48), Inches(2.65), Inches(2.75),
                size=11, color=C_TEXT)

    # 矢印
    txb(s, "→", Inches(3.1), Inches(2.95), Inches(0.55), Inches(0.4),
        size=20, bold=True, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
    txb(s, "→", Inches(6.15), Inches(2.95), Inches(0.55), Inches(0.4),
        size=20, bold=True, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

    rect(s, M, Inches(5.45), IW, Inches(0.58), C_CODE_BG, border=True)
    rect(s, M, Inches(5.45), Inches(0.055), Inches(0.58), C_TEAL)
    txb(s, "公開されるのは within30 フラグ と Poseidon commitment ハッシュのみ — 生座標は外へ出ない",
        M + Inches(0.15), Inches(5.54), IW - Inches(0.2), Inches(0.4),
        size=12, bold=True, color=C_TEAL)

    # ── スライド 3: システムアーキテクチャ ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "システムアーキテクチャ",
           "ハードウェア → BLE → Flutter アプリ → オンデバイス ZK 証明 の全体フロー")
    page_num(s, 3)

    # レイヤー図（行ごとに描画）
    layers = [
        (C_NAVY,  "① Raspberry Pi (UAV 搭載)",
         "Python remote_id.py で BasicID / Location を ASTM F3411-22a 形式で BLE ブロードキャスト（1秒ごと）"),
        (C_BLUE,  "② BLE 伝送 (OpenDroneID)",
         "Service UUID 0xFFFA の Service Data AD / 25 バイトメッセージ / 距離 ~100m 圏内に届く"),
        (C_TEAL,  "③ Flutter アプリ (Android / Pixel 8)",
         "BLE スキャン → Location パース → GPS 取得 → カメラ録画 → 履歴 20 点収集"),
        (RGBColor(0x6A, 0x1B, 0x9A), "④ Circom 回路 (オンデバイス)",
         "Witness 計算 (circom_witnesscalc/.wcd) → Groth16 証明生成 (flutter_rapidsnark/.zkey) → 検証"),
        (RGBColor(0xC6, 0x28, 0x28), "⑤ Public Outputs (証明結果)",
         "within30 / any_within30 / drone_commitment / target_commitment — 生座標ゼロ公開"),
    ]
    for i, (col, lbl, desc) in enumerate(layers):
        top = Inches(1.15) + Inches(1.1 * i)
        rect(s, M, top, Inches(2.3), Inches(0.85), col)
        txb(s, lbl, M + Inches(0.1), top + Inches(0.22), Inches(2.1), Inches(0.42),
            size=10, bold=True, color=C_WHITE)
        rect(s, M + Inches(2.3), top, IW - Inches(2.3), Inches(0.85), C_CODE_BG, border=True)
        txb(s, desc, M + Inches(2.45), top + Inches(0.2), IW - Inches(2.6), Inches(0.5),
            size=11, color=C_TEXT)
        if i < len(layers) - 1:
            txb(s, "▼", M + Inches(0.8), top + Inches(0.85), Inches(0.7), Inches(0.25),
                size=9, color=col, align=PP_ALIGN.CENTER)

    # ── スライド 4: 開発環境・ツールバージョン ───────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "開発環境・ツールバージョン",
           "ビルドおよびデプロイに必要なすべてのツール")
    page_num(s, 4)

    table(s,
          ["ツール", "バージョン", "役割", "インストール方法"],
          [
              ["circom",         "2.2.2",          "Circom 回路コンパイラ (.r1cs / .wasm 生成)",
               "cargo install circom  (Rust)"],
              ["snarkjs",        "0.7.5",          "Powers of Tau / Groth16 セットアップ / 証明検証",
               "npm install -g snarkjs"],
              ["Node.js",        "v24.14.0",       "snarkjs の実行環境",
               "https://nodejs.org"],
              ["circomlib",      "^2.0.5 (npm)",   "Poseidon ハッシュ等の標準回路ライブラリ",
               "npm install  (package.json 記載)"],
              ["build-circuit",  "0.1.1 (git)",    ".wcd ファイル生成 (circom-witnesscalc)",
               "cargo build -p build-circuit\n(iden3/circom-witnesscalc)"],
              ["Flutter",        "3.32.5 stable",  "Android アプリビルド",
               "https://flutter.dev"],
              ["Dart SDK",       "^3.8.1",         "Flutter / Dart 言語ランタイム",
               "Flutter に同梱"],
              ["Python",         "3.14.4",         "remote_id.py (Raspberry Pi BLE 送信)",
               "OS 標準 / pyenv"],
              ["rustc",          "1.90.0",         "circom / build-circuit のビルド",
               "https://rustup.rs"],
          ],
          M, Inches(1.12), IW,
          col_widths=[Inches(1.5), Inches(1.4), Inches(3.6), Inches(2.6)],
          font_size=10)

    # ── スライド 5: ZK 回路ツールチェーン ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "ZK 回路ツールチェーン",
           "回路の作成からアプリ搭載まで — scripts/build_*.sh で一括実行")
    page_num(s, 5)

    # 左: フロー
    txb(s, "ビルドフロー",
        M, Inches(1.12), Inches(4.1), Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    steps = [
        ("circom --O2",        "回路コンパイル → .r1cs / .wasm 生成"),
        ("snarkjs powersoftau", "Powers of Tau (bn128 / pot12~14)"),
        ("snarkjs groth16 setup","Groth16 セットアップ → _0000.zkey"),
        ("snarkjs zkey contribute","Zkey contribution → _final.zkey"),
        ("snarkjs zkey export",  "検証キー → verification_key.json"),
        ("snarkjs wtns calculate","サンプル Witness 生成"),
        ("snarkjs groth16 prove", "サンプル Proof 生成・確認"),
        ("build-circuit",       ".wcd ファイル生成 (circom-witnesscalc)"),
    ]
    for i, (cmd, desc) in enumerate(steps):
        top = Inches(1.45) + Inches(0.62 * i)
        rect(s, M, top, Inches(0.22), Inches(0.44), C_BLUE)
        txb(s, str(i + 1), M + Inches(0.02), top + Inches(0.1), Inches(0.18), Inches(0.28),
            size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txb(s, cmd, M + Inches(0.28), top + Inches(0.02), Inches(1.8), Inches(0.22),
            size=9, bold=True, color=C_NAVY)
        txb(s, desc, M + Inches(0.28), top + Inches(0.22), Inches(3.6), Inches(0.2),
            size=9, color=C_SUBTEXT)

    # 右: ライブラリ詳細
    txb(s, "主要ライブラリ・リポジトリ",
        Inches(5.1), Inches(1.12), Inches(4.65), Inches(0.3),
        size=11, bold=True, color=C_NAVY)

    zk_libs = [
        ("circom", "github.com/iden3/circom",
         "Circom 言語コンパイラ (Rust製)\n.circom → .r1cs + .wasm"),
        ("snarkjs", "github.com/iden3/snarkjs",
         "ZK-SNARK ツールキット (JS/Node)\nGroth16 / PLONK / Powers of Tau"),
        ("circomlib", "github.com/iden3/circomlib",
         "標準回路ライブラリ (npm)\nPoseidon / EdDSA / Comparators など"),
        ("circom-witnesscalc", "github.com/iden3/circom-witnesscalc",
         "Rust 製 Witness 計算エンジン\nbuild-circuit で .wcd を生成"),
    ]
    for i, (name, repo, desc) in enumerate(zk_libs):
        top = Inches(1.45) + Inches(1.35 * i)
        rect(s, Inches(5.1), top, Inches(4.55), Inches(1.2), C_CODE_BG, border=True)
        rect(s, Inches(5.1), top, Inches(0.055), Inches(1.2), C_BLUE)
        txb(s, name, Inches(5.22), top + Inches(0.08), Inches(4.2), Inches(0.28),
            size=11, bold=True, color=C_NAVY)
        txb(s, repo, Inches(5.22), top + Inches(0.35), Inches(4.2), Inches(0.22),
            size=9, color=C_BLUE, italic=True)
        txb(s, desc, Inches(5.22), top + Inches(0.57), Inches(4.2), Inches(0.55),
            size=10, color=C_TEXT)

    # ── スライド 6: Flutter アプリ依存ライブラリ ────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "Flutter アプリ依存ライブラリ",
           "rapidsnark_app/pubspec.yaml — pub.dev / GitHub から取得")
    page_num(s, 6)

    # 左カラム
    txb(s, "ZK 証明コア",
        M, Inches(1.12), Inches(4.3), Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["パッケージ", "バージョン", "用途"],
          [
              ["flutter_rapidsnark", "git (iden3)", "Groth16 証明生成・検証 (rapidsnark Rust バインディング)"],
              ["circom_witnesscalc", "^0.0.1-alpha.1", ".wcd を使ったオンデバイス Witness 計算"],
          ],
          M, Inches(1.42), Inches(4.3),
          col_widths=[Inches(1.8), Inches(0.95), Inches(1.55)], font_size=10)

    txb(s, "センサー・通信",
        M, Inches(2.75), Inches(4.3), Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["パッケージ", "バージョン", "用途"],
          [
              ["flutter_blue_plus", "^1.35.5", "BLE スキャン (OpenDroneID 受信)"],
              ["geolocator",        "^14.0.2",  "GPS ストリーム (1 秒ごと位置取得)"],
              ["sensors_plus",      "^7.0.0",   "気圧センサー → 高度推定"],
              ["camera",            "^0.11.2+1","カメラ録画 (動画証拠)"],
          ],
          M, Inches(3.05), Inches(4.3),
          col_widths=[Inches(1.8), Inches(0.95), Inches(1.55)], font_size=10)

    # 右カラム
    txb(s, "AI / 映像解析",
        Inches(5.1), Inches(1.12), Inches(4.55), Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["パッケージ", "バージョン", "用途"],
          [
              ["google_mlkit_image_labeling",   "^0.14.2", "ML Kit 画像ラベリング (UAV 検知)"],
              ["google_mlkit_object_detection", "^0.15.1", "ML Kit 物体検出 (バウンディングボックス)"],
              ["tflite_flutter",                "^0.11.0", "TFLite モデル推論 (カスタム UAV モデル)"],
              ["video_thumbnail",               "^0.5.6",  "動画サムネイル / フレーム抽出"],
              ["image",                         "^4.5.4",  "画像デコード・前処理"],
          ],
          Inches(5.1), Inches(1.42), Inches(4.55),
          col_widths=[Inches(2.15), Inches(0.85), Inches(1.55)], font_size=10)

    txb(s, "ユーティリティ",
        Inches(5.1), Inches(3.62), Inches(4.55), Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["パッケージ", "バージョン", "用途"],
          [
              ["path_provider", "^2.1.5", "一時ファイルパス (.zkey コピー先)"],
              ["image_picker",  "^1.2.0", "ギャラリーから動画選択"],
          ],
          Inches(5.1), Inches(3.92), Inches(4.55),
          col_widths=[Inches(1.5), Inches(0.85), Inches(2.2)], font_size=10)

    # ── スライド 7: リポジトリ構成 ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "リポジトリ構成",
           "github.com/SMRI2170/zkp-m1")
    page_num(s, 7)

    txb(s, "github.com/SMRI2170/zkp-m1",
        M, Inches(1.1), IW, Inches(0.3),
        size=13, bold=True, color=C_BLUE)

    code(s,
         "zkp-m1/\n"
         "├── circuits/                   # Circom 回路ソース\n"
         "│   ├── lib_uav_safety.circom   #   共通テンプレート (Distance30mTimeCore)\n"
         "│   ├── distance_30m_time.circom\n"
         "│   ├── distance_30m_time_detection_flag.circom\n"
         "│   ├── distance_30m_time_detection_commitment.circom\n"
         "│   ├── distance_30m_time_feature_classifier.circom\n"
         "│   └── distance_30m_time_series.circom  # N=20 時系列回路\n"
         "├── scripts/                    # ビルドスクリプト群\n"
         "│   ├── build_distance_variant_assets.sh   # 共通ビルドテンプレート\n"
         "│   └── build_distance_*_assets.sh         # 各回路ラッパー\n"
         "├── zk-distance/\n"
         "│   └── samples/                # サンプル入力 JSON (各回路)\n"
         "├── rapidsnark_app/             # Flutter アプリ (Android)\n"
         "│   ├── lib/main.dart           #   アプリ本体\n"
         "│   └── assets/circuits/        #   .wcd / .zkey / verification_key.json\n"
         "├── uav/\n"
         "│   └── remote_id.py            # Raspberry Pi BLE 送信スクリプト\n"
         "└── package.json                # circomlib npm 依存",
         M, Inches(1.45), IW, Inches(5.6), size=10)

    # ── スライド 8: OpenDroneID とは ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "1.  OpenDroneID とは",
           "国際標準 ASTM F3411 — ドローンが BLE/Wi-Fi で識別・位置情報をブロードキャストする規格")
    page_num(s, 2)

    bullets(s, [
        "■  ドローンが BLE または Wi-Fi で自身の識別情報・位置情報を周囲に送信する国際規格",
        "■  受信側は特別なアプリ不要で機体識別・位置を把握できる「電波リモートID」",
        "■  日本: 2022年6月 改正航空法により 100g 以上の機体に搭載義務",
    ], M, Inches(1.15), IW, Inches(1.1), size=12)

    table(s,
          ["規格名", "内容"],
          [
              ["ASTM F3411-22a", "米国 ASTM International 制定（2022年版）"],
              ["ASD-STAN prEN 4709-002", "欧州版 — ASTM と相互運用性あり"],
              ["日本航空法 改正（2022年6月）", "100g 以上の機体にリモートID搭載義務"],
          ],
          M, Inches(2.35), IW,
          col_widths=[Inches(3.4), Inches(5.7)])

    rect(s, M, Inches(4.35), IW, Inches(0.52), C_CODE_BG, border=True)
    rect(s, M, Inches(4.35), Inches(0.055), Inches(0.52), C_TEAL)
    txb(s, "キーポイント: 地上・上空にいる第三者が機体識別・現在位置をリアルタイムで受信できる",
        M + Inches(0.15), Inches(4.42), IW - Inches(0.2), Inches(0.38),
        size=12, bold=True, color=C_TEAL)

    # ── スライド 3: 変更の概要 ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "2.  変更の概要",
           "独自 BLE フォーマット → ASTM F3411-22a (OpenDroneID) 準拠")
    page_num(s, 3)

    txb(s, "変更前（独自フォーマット）",
        M, Inches(1.12), Inches(4.3), Inches(0.32),
        size=11, bold=True, color=C_SUBTEXT)
    code(s,
         "BLE Manufacturer Specific Data\n"
         "  AD type   : 0xFF\n"
         "  Company ID: 0x02E5\n"
         "  Payload   : [MAGIC=0x5A][VER=0x01]\n"
         "              [lat_e7 4B][lon_e7 4B]\n"
         "              [alt_dm 2B][unix_time 4B][seq 1B]\n"
         "  合計      : 17 バイト（独自仕様）",
         M, Inches(1.46), Inches(4.3), Inches(2.3), size=10)

    txb(s, "→", Inches(4.85), Inches(2.3), Inches(0.55), Inches(0.5),
        size=28, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    txb(s, "変更後（OpenDroneID 準拠）",
        Inches(5.5), Inches(1.12), Inches(4.3), Inches(0.32),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "BLE Service Data\n"
         "  AD type     : 0x16\n"
         "  Service UUID: 0xFFFA\n"
         "  Payload     : 25 バイトの ASTM F3411\n"
         "                メッセージ\n"
         "  合計        : 29 バイト（規格準拠）\n"
         "                BLE 31B 制限内に収まる",
         Inches(5.5), Inches(1.46), Inches(4.3), Inches(2.3), size=10)

    rect(s, M, Inches(3.95), IW, Pt(1), C_RULE)
    txb(s, "得られるメリット",
        M, Inches(4.08), IW, Inches(0.32),
        size=12, bold=True, color=C_NAVY)
    bullets(s, [
        "■  DJI / Autel などの市販 Remote ID モジュールと互換",
        "■  DroneScanner / OpenDroneID Receiver などの対応アプリで受信可能",
        "■  規制当局の検査機器（FAA・国交省）に対応",
    ], M, Inches(4.42), IW, Inches(1.2), size=12)

    # ── スライド 4: BLE ペイロード構造 ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "3.  BLE ペイロード構造",
           "AD 要素レイアウト（合計 29 バイト）— BLE 31 バイト広告制限内に収まる")
    page_num(s, 4)

    code(s,
         "オクテット:  [0]      [1]    [2]   [3]    [4 .. 28]\n"
         "内容:        length   0x16   0xFA  0xFF   25 バイトメッセージ\n"
         "             (= 28)   AD type  UUID (Little-Endian)\n\n"
         "  length=28  : type(1) + UUID(2) + message(25) の合計バイト数\n"
         "  AD type 0x16 : Service Data – 16-bit UUID\n"
         "  UUID 0xFFFA  : OpenDroneID 専用 Service UUID",
         M, Inches(1.12), IW, Inches(2.15), size=11)

    txb(s, "BasicID メッセージ  (25 B)",
        M, Inches(3.42), Inches(4.4), Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "Byte 0    : Header = (Type=0 << 4) | ProtoVer\n"
         "Byte 1    : IDType[3:0] | UAType[7:4]\n"
         "Bytes 2-21: UAS_ID  (20B ASCII, null padding)\n"
         "Bytes 22-24: Reserved  (0x00)",
         M, Inches(3.74), Inches(4.4), Inches(1.7), size=10)

    txb(s, "Location メッセージ  (25 B)",
        Inches(5.1), Inches(3.42), Inches(4.4), Inches(0.3),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "Byte 0    : Header = (Type=1 << 4) | ProtoVer\n"
         "Bytes 5-8 : Latitude   (int32 LE, 1e-7°)\n"
         "Bytes 9-12: Longitude  (int32 LE, 1e-7°)\n"
         "Bytes 15-16: AltGeo   (uint16, (alt+1000)×2)\n"
         "Bytes 21-22: Timestamp (UTC 時間内秒 × 10)",
         Inches(5.1), Inches(3.74), Inches(4.4), Inches(1.7), size=10)

    # ── スライド 5: BasicID 詳細 ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "3-1.  BasicID メッセージ詳細",
           "機体識別子と種別 — 25 バイト固定長")
    page_num(s, 5)

    table(s,
          ["バイト", "フィールド", "値・説明"],
          [
              ["Byte 0", "Header",
               "(MessageType=0 << 4) | ProtoVersion　例: 0x02"],
              ["Byte 1", "IDType | UAType",
               "IDType=1(Serial) | UAType=2(Multirotor)<<4　例: 0x21"],
              ["Bytes 2–21", "UAS_ID",
               '20 バイト ASCII / null パディング　例: "UAV-ZKP-M1-DEMO00"'],
              ["Bytes 22–24", "Reserved", "0x00  0x00  0x00"],
          ],
          M, Inches(1.12), IW,
          col_widths=[Inches(1.5), Inches(2.3), Inches(5.3)])

    txb(s, "実際の送信例（hex）",
        M, Inches(3.42), IW, Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "02 21 55 41 56 2D 5A 4B 50 2D 4D 31 2D 44 45 4D 4F 30 30 00 00 00 00 00 00\n"
         "↑   ↑  ↑──────── \"UAV-ZKP-M1-DEMO00\" ──────────────────────────↑  ↑── reserved\n"
         "型  ID UAV種別",
         M, Inches(3.75), IW, Inches(1.5), size=10)

    # ── スライド 6: Location 詳細 ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "3-2.  Location メッセージ詳細",
           "現在位置・高度・速度・タイムスタンプ — 1 秒に 1 回以上の送信が義務")
    page_num(s, 6)

    table(s,
          ["バイト", "フィールド", "エンコード"],
          [
              ["Byte 0",      "Header",       "(Type=1<<4)|ProtoVer　例: 0x12"],
              ["Byte 1",      "Status | Flags","Status=1(Airborne)　例: 0x01"],
              ["Bytes 5–8",   "Latitude",     "int32 LE, 1e-7 度　35.6891390° → 356,891,390"],
              ["Bytes 9–12",  "Longitude",    "int32 LE, 1e-7 度　139.6917000° → 1,396,917,000"],
              ["Bytes 13–14", "AltitudeBaro", "uint16 LE, (alt_m + 1000) × 2"],
              ["Bytes 15–16", "AltitudeGeo",  "uint16 LE, (alt_m + 1000) × 2"],
              ["Bytes 21–22", "Timestamp",    "uint16 LE, UTC 時間内秒 × 10　例: 19200"],
          ],
          M, Inches(1.12), IW,
          col_widths=[Inches(1.5), Inches(2.1), Inches(5.5)])

    txb(s, "実際の送信例",
        M, Inches(4.95), IW, Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "12 01 00 00 00 FE BA 45 15 08 43 43 53 E8 07 E8 07 FF FF 00 00 6E 4B 00 00\n"
         "↑  ↑          ↑──────────↑ ↑──────────↑ baro geo                ↑───↑ timestamp",
         M, Inches(5.28), IW, Inches(1.1), size=10)

    # ── スライド 7: 送信側 (Raspberry Pi) ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "4.  送信側実装（Raspberry Pi）",
           "uav/remote_id.py — hcitool で Raw HCI コマンドを直接発行")
    page_num(s, 7)

    txb(s, "送信サイクル（1 秒ごと交互）",
        M, Inches(1.12), Inches(4.2), Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "tick 0: Location  ← 現在位置・高度・時刻\n"
         "tick 1: BasicID   ← 機体識別子・種別\n"
         "tick 2: Location\n"
         "tick 3: BasicID  ...\n"
         "→ 8 秒スキャンで Location を最低 4 回受信可",
         M, Inches(1.44), Inches(4.2), Inches(1.95), size=10)

    txb(s, "HCI コマンド送信フロー",
        Inches(5.0), Inches(1.12), Inches(4.65), Inches(0.3),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "# 広告パラメータ設定 (ADV_NONCONN_IND)\n"
         "hcitool cmd 0x08 0x0006 A0 00 A0 00 03 ...\n\n"
         "# 広告データ設定 (29 バイト)\n"
         "hcitool cmd 0x08 0x0008\n"
         "  1D 1C 16 FA FF [25B msg] 00 00\n\n"
         "# 広告有効化\n"
         "hcitool cmd 0x08 0x000A 01",
         Inches(5.0), Inches(1.44), Inches(4.65), Inches(1.95), size=10)

    rect(s, M, Inches(3.52), IW, Pt(1), C_RULE)
    txb(s, "使い方",
        M, Inches(3.65), IW, Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "# 固定座標を 1 秒ごとに送信\n"
         "python3 uav/remote_id.py --lat 35.6891390 --lon 139.6917000 --alt 12.0\n\n"
         "# JSON ファイルから動的に読み込む（飛行中に更新可能）\n"
         "python3 uav/remote_id.py --json-path uav_sample.json --interval 0.5",
         M, Inches(3.98), IW, Inches(1.65), size=10)

    # ── スライド 8: 受信側 (Flutter) ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "5.  受信側実装（Flutter アプリ）",
           "rapidsnark_app/lib/main.dart — UavTelemetryPacket クラス")
    page_num(s, 8)

    txb(s, "BLE スキャンフロー",
        M, Inches(1.12), Inches(3.5), Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    bullets(s, [
        "①  Bluetooth アダプタ状態を確認",
        "②  FlutterBluePlus.startScan (25 秒)",
        "③  serviceData から UUID 0xFFFA を検索",
        "④  BasicID → UAS_ID を保存",
        "⑤  Location → UavTelemetryPacket に変換",
        "⑥  フォームに lat / lon / alt / time を自動入力",
    ], M, Inches(1.44), Inches(3.5), Inches(2.8), size=11)

    txb(s, "Location パース & タイムスタンプ再構築",
        Inches(4.2), Inches(1.12), Inches(5.5), Inches(0.3),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "// Service UUID 0xFFFA のデータを取得\n"
         "static const serviceUuid = '0000fffa-...';\n\n"
         "// Location メッセージ (type=1) → パース\n"
         "final latI32 = data.getInt32(5,  Endian.little);\n"
         "final lonI32 = data.getInt32(9,  Endian.little);\n"
         "final altGeo = data.getUint16(15, Endian.little);\n"
         "final tsRaw  = data.getUint16(21, Endian.little);\n\n"
         "// タイムスタンプ: 時間内秒 → Unix 秒\n"
         "final secs = tsRaw ~/ 10;\n"
         "final base = nowSec - (nowSec % 3600);\n"
         "return base + secs;",
         Inches(4.2), Inches(1.44), Inches(5.5), Inches(4.0), size=10)

    # ── スライド 9: フォーマット対応表 ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "6.  独自フォーマット vs OpenDroneID",
           "フィールド対応表")
    page_num(s, 9)

    table(s,
          ["フィールド", "独自フォーマット", "OpenDroneID (ASTM F3411)"],
          [
              ["識別子",         "Manufacturer Data (0x02E5)",          "Service Data, UUID 0xFFFA"],
              ["緯度・経度",      "int32 LE, 1e-7 度",                   "同じ"],
              ["高度",            "int16 LE, デシメートル (0.1m単位)",    "uint16 LE, (alt+1000)×2 (0.5m単位)"],
              ["タイムスタンプ",  "uint32 LE, Unix 秒（絶対値）",         "uint16 LE, 時間内秒×10（相対値）"],
              ["機体識別子",      "なし",                                 "BasicID の UAS_ID（20B ASCII）"],
              ["速度・方向",      "なし",                                 "Direction / SpeedH / SpeedV"],
              ["精度情報",        "なし",                                 "HorizAcc / VertAcc / SpeedAcc など"],
          ],
          M, Inches(1.12), IW,
          col_widths=[Inches(1.8), Inches(3.1), Inches(4.2)])

    # ── スライド 10: ZK 証明フロー ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "8.  ZK 証明フローへの接続",
           "OpenDroneID Location → Circom 回路 private input への変換")
    page_num(s, 10)

    code(s,
         "Location.latitude     →  drone_lat_e7_norm  =  round((lat + 90)  × 1e7)\n"
         "Location.longitude    →  drone_lon_e7_norm  =  round((lon + 180) × 1e7)\n"
         "Location.altitudeGeo  →  drone_alt_cm       =  round(alt_m × 100)\n"
         "Location.timestamp    →  drone_time_sec     =  _reconstructUnixTime(tsRaw)",
         M, Inches(1.12), IW, Inches(1.65), size=11)

    txb(s, "回路内での処理",
        M, Inches(2.9), IW, Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    bullets(s, [
        "■  Poseidon4(drone_lat, drone_lon, drone_alt, drone_time) → drone_commitment  【公開】",
        "■  Poseidon4(target_lat, target_lon, target_alt, target_time) → target_commitment  【公開】",
        "■  実際の座標値は外部に一切公開されない — ハッシュ値のみが public output",
        "■  within30 = (3D距離 ≤ 30m) AND (時刻差 ≤ 5s) → 公開フラグとして証明",
    ], M, Inches(3.24), IW, Inches(1.85), size=12)

    rect(s, M, Inches(5.18), IW, Inches(0.62), C_CODE_BG, border=True)
    rect(s, M, Inches(5.18), Inches(0.055), Inches(0.62), C_TEAL)
    txb(s, "Time Series 回路 (N=20): 録画中に収集した 20 ペア全点を対象に OR 削減 → any_within30",
        M + Inches(0.15), Inches(5.27), IW - Inches(0.2), Inches(0.42),
        size=12, bold=True, color=C_TEAL)

    # ── スライド 11: 今後の拡張方針 ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "9.  今後の拡張方針",
           "ロードマップ")
    page_num(s, 11)

    table(s,
          ["項目", "内容"],
          [
              ["BT5 Long Range",
               "Extended Advertising + MessagePack で複数メッセージを 1 パケットで送信"],
              ["Wi-Fi NaN",
               "より広範囲への送信（Pixel 8 は Wi-Fi NaN 受信対応）"],
              ["Authentication",
               "OpenDroneID Auth メッセージで機体の署名を検証"],
              ["IPFS 連携",
               "BasicID の UAS_ID をキーに IPFS から Merkle Tree を取得"],
              ["Merkle 証明",
               "飛行ログの Merkle Proof を ZK 回路に組み込み改ざん防止"],
          ],
          M, Inches(1.12), IW,
          col_widths=[Inches(2.2), Inches(6.9)])

    rect(s, M, Inches(4.4), IW, Inches(0.58), C_CODE_BG, border=True)
    rect(s, M, Inches(4.4), Inches(0.055), Inches(0.58), C_NAVY)
    txb(s, "目標: OpenDroneID 準拠 + ZK 証明 + 飛行ログ Merkle Proof の統合アーキテクチャ",
        M + Inches(0.15), Inches(4.49), IW - Inches(0.2), Inches(0.4),
        size=12, bold=True, color=C_NAVY)

    # ── スライド 12: GPS + BME280 実機センサー統合 ──────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "GPS + BME280 実機センサー統合",
           "Raspberry Pi に接続した実機センサーからリアルタイムでドローン位置・高度を取得")
    page_num(s, 12)

    # ハードウェア構成
    txb(s, "ハードウェア構成",
        M, Inches(1.12), IW, Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["デバイス", "接続方式", "識別子", "役割"],
          [
              ["GPS モジュール（u-blox）", "USB CDC", "/dev/ttyACM0", "緯度・経度・GPS 高度・Fix 状態"],
              ["BME280 気圧センサー",       "I²C",    "アドレス 0x76", "気圧 → QFE 相対高度計算"],
              ["Raspberry Pi",             "—",      "192.168.11.54", "BLE ブロードキャスト・センサー統合"],
          ],
          M, Inches(1.44), IW,
          col_widths=[Inches(2.1), Inches(1.5), Inches(1.8), Inches(3.7)], font_size=11)

    # 変更前後
    txb(s, "変更前（固定値）",
        M, Inches(2.98), Inches(4.3), Inches(0.28),
        size=11, bold=True, color=C_SUBTEXT)
    code(s,
         "python3 remote_id.py \\\n"
         "  --lat 35.6891390 \\\n"
         "  --lon 139.6917000 \\\n"
         "  --alt 12.0\n"
         "# 座標はCLI引数の固定値",
         M, Inches(3.28), Inches(4.3), Inches(1.72), size=10)

    txb(s, "→", Inches(4.85), Inches(3.9), Inches(0.55), Inches(0.5),
        size=28, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    txb(s, "変更後（実機センサー）",
        Inches(5.5), Inches(2.98), Inches(4.15), Inches(0.28),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "python3 remote_id.py --gps\n"
         "# GPS シリアル自動検出\n"
         "# BME280 I2C 読み取り\n"
         "# 1秒ごとにリアルタイム送信",
         Inches(5.5), Inches(3.28), Inches(4.15), Inches(1.72), size=10)

    # 実行結果
    rect(s, M, Inches(5.08), IW, Inches(0.64), C_CODE_BG, border=True)
    rect(s, M, Inches(5.08), Inches(0.055), Inches(0.64), C_TEAL)
    txb(s, "実機確認結果: GPS Fix 取得、BME280 正常、BLE ブロードキャスト確認\n"
        "[seq=001] lat=34.9686023  lon=135.8991972  alt_geo=122.1m  alt_baro=0.1m",
        M + Inches(0.15), Inches(5.15), IW - Inches(0.2), Inches(0.52),
        size=10, color=C_TEAL)

    # ── スライド 13: QFE 高度方式 ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "QFE 高度方式 — 同一基準での高度計測",
           "ドローン（Pi BME280）とユーザー（スマホ気圧計）を同じ地上基準ゼロに揃える")
    page_num(s, 13)

    # 問題
    txb(s, "問題: なぜ高度がずれるか",
        M, Inches(1.12), IW, Inches(0.3),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["方式", "問題点"],
          [
              ["標準大気 1013.25 hPa 基準",
               "実際の海面気圧は天気で変化（±10 hPa 以上）→ 絶対高度が最大 ±80m ずれる"],
              ["GPS 高度 (WGS84)",
               "楕円体高で誤差が大きい（±10〜30m 以上）、室内・屋根下では取得不可"],
              ["Pi と スマホで基準が違う",
               "BME280 と スマホ気圧計が異なる P₀ を使うと高度差が正確に出ない"],
          ],
          M, Inches(1.44), IW,
          col_widths=[Inches(2.5), Inches(6.6)], font_size=11)

    # 解決策
    txb(s, "解決策: QFE（Height Above Takeoff）",
        M, Inches(3.32), IW, Inches(0.3),
        size=11, bold=True, color=C_TEAL)

    # 2カラム
    code(s,
         "# Pi 側（remote_id.py）\n"
         "# 起動時に地上気圧を記録\n"
         "P0 = bme280.sample().pressure\n\n"
         "# 相対高度を計算\n"
         "alt = 44330 × (1 - (P/P0)^(1/5.255))\n\n"
         "# 地上では alt ≈ 0m\n"
         "# 10m 飛行後 alt ≈ 10m",
         M, Inches(3.62), Inches(4.3), Inches(2.7), size=10)

    code(s,
         "// Flutter 側（main.dart）\n"
         "// 最初の気圧値を地上基準に\n"
         "_groundPressureHpa ??= event.pressure;\n\n"
         "// QFE 高度計算\n"
         "double _baroToRelativeAlt(double p) {\n"
         "  return 44330 * (1 - (p/p0)^(1/5.255));\n"
         "}\n"
         "// スマホも地上 = 0m",
         Inches(5.5), Inches(3.62), Inches(4.15), Inches(2.7), size=10)

    # 効果の説明
    rect(s, M, Inches(6.42), IW, Inches(0.7), C_CODE_BG, border=True)
    rect(s, M, Inches(6.42), Inches(0.055), Inches(0.7), C_BLUE)
    txb(s, "効果: Pi と スマホが同じ場所で起動 → P₀ がほぼ同値 → 高度差 = 実際の物理的な高さの差\n"
        "天気・標高・センサーの個体差に関わらず、両者間の相対高度が正確になる",
        M + Inches(0.15), Inches(6.5), IW - Inches(0.2), Inches(0.55),
        size=11, bold=True, color=C_BLUE)

    # ── スライド 14: Flutter アプリ改善 ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "Flutter アプリ改善",
           "録画安定化・GPS 1秒更新・高度基準統一の 3 点を修正")
    page_num(s, 14)

    for i, (icon, title, before, after, reason) in enumerate([
        ("①", "録画エラーの可視化",
         "catch (_) {}\n  // 全エラーを無視",
         "catch (e) {\n  setState(() =>\n    _status = '録画エラー: $e');\n}",
         "無音でのエラーを表示に変更し、原因特定を容易にする"),
        ("②", "imageFormatGroup 削除",
         "CameraController(\n  imageFormatGroup:\n    ImageFormatGroup.jpeg,\n)",
         "CameraController(\n  // imageFormatGroup なし\n  enableAudio: false,\n)",
         "JPEG フォーマット指定が動画録画と競合するため削除"),
        ("③", "GPS 1 秒更新",
         "LocationSettings(\n  distanceFilter: 0,\n  // 移動しないと更新なし\n)",
         "AndroidSettings(\n  intervalDuration:\n    Duration(seconds: 1),\n)",
         "停止中でも毎秒 GPS を取得、時系列データの品質向上"),
    ]):
        top = Inches(1.12) + Inches(1.98 * i)
        txb(s, f"{icon}  {title}",
            M, top, IW, Inches(0.3),
            size=11, bold=True, color=C_NAVY)
        rect(s, M, top, Inches(0.055), Inches(0.3), C_BLUE)
        code(s, before, M, top + Inches(0.32), Inches(3.8), Inches(1.42), size=10)
        txb(s, "→", Inches(4.1), top + Inches(0.75), Inches(0.55), Inches(0.4),
            size=20, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        code(s, after, Inches(4.7), top + Inches(0.32), Inches(3.8), Inches(1.42), size=10)
        txb(s, reason, M, top + Inches(1.76), IW, Inches(0.22),
            size=9, color=C_SUBTEXT, italic=True)

    # ── スライド 16: opendroneid-core-c 準拠 — 送信側の誤り修正 ────────────────
    s = prs.slides.add_slide(blank)
    header(s, "opendroneid-core-c 準拠 — 送信側の誤り修正",
           "remote_id.py に存在した 3 箇所のビットレイアウト不整合を修正")
    page_num(s, 16)

    txb(s, "修正前（誤）",
        M, Inches(1.12), Inches(4.3), Inches(0.28),
        size=11, bold=True, color=RGBColor(0xC6, 0x28, 0x28))
    code(s,
         "# ❌ Status が下位ニブル（正しくは上位）\n"
         "byte1 = (status & 0x0F) | (flags << 4)\n\n"
         "# ❌ AIRBORNE = 1（正しくは 2）\n"
         "ODID_STATUS_AIRBORNE = 1\n\n"
         "# ❌ BasicID Byte1 のニブルが逆\n"
         "byte1 = (ua_type << 4) | id_type",
         M, Inches(1.42), Inches(4.3), Inches(2.15), size=10)

    txb(s, "→", Inches(4.85), Inches(2.2), Inches(0.55), Inches(0.5),
        size=28, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    txb(s, "修正後（opendroneid-core-c 準拠）",
        Inches(5.5), Inches(1.12), Inches(4.15), Inches(0.28),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "# ✓ Status が上位ニブル (bits 7:4)\n"
         "byte1 = ((status & 0x0F) << 4) | flags\n\n"
         "# ✓ AIRBORNE = 2 (GROUND=1, AIRBORNE=2)\n"
         "ODID_STATUS_AIRBORNE = 2\n\n"
         "# ✓ BasicID: IDType 上位, UAType 下位\n"
         "byte1 = ((id_type & 0x0F) << 4) | ua_type",
         Inches(5.5), Inches(1.42), Inches(4.15), Inches(2.15), size=10)

    txb(s, "Location Byte 1 の正しいビットレイアウト（opendroneid-core-c より）",
        M, Inches(3.68), IW, Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "Byte 1: [ Status[7:4] | Reserved[3] | HeightType[2] | EWDir[1] | SpeedMult[0] ]\n"
         "         ^^^^^^^^^^^^\n"
         "         上位 4 ビット — UNDECLARED=0, GROUND=1, AIRBORNE=2, EMERGENCY=3",
         M, Inches(3.98), IW, Inches(1.48), size=10)

    rect(s, M, Inches(5.55), IW, Inches(0.44), C_CODE_BG, border=True)
    rect(s, M, Inches(5.55), Inches(0.055), Inches(0.44), C_TEAL)
    txb(s, "方向・速度のエンコードも追加: EWDirection フラグで西半球補正 / SpeedMult で高速域スケール切替",
        M + Inches(0.15), Inches(5.62), IW - Inches(0.2), Inches(0.3),
        size=11, color=C_TEAL)

    # ── スライド 17: Flutter 受信側デコーダの更新 ────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "Flutter 受信側デコーダの更新",
           "UavTelemetryPacket を opendroneid-core-c 準拠に書き直し")
    page_num(s, 17)

    txb(s, "Byte 1 の正しいデコード",
        M, Inches(1.12), Inches(4.3), Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "// 修正前: Status を下位ニブルから読んでいた\n"
         "// (byte1 & 0x0F) → 誤り\n\n"
         "// 修正後: Status を上位ニブルから読む\n"
         "final byte1     = rawBytes[1] & 0xFF;\n"
         "final status    = (byte1 >> 4) & 0x0F;  // bits 7:4\n"
         "final ewDir     = (byte1 >> 1) & 0x01;  // bit 1\n"
         "final speedMult = byte1 & 0x01;          // bit 0\n\n"
         "statusIsAirborne: status == 2,  // AIRBORNE=2",
         M, Inches(1.42), Inches(4.3), Inches(3.05), size=10)

    txb(s, "方向・速度デコードの追加",
        Inches(5.5), Inches(1.12), Inches(4.15), Inches(0.28),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "// EWDirection フラグで西半球補正\n"
         "double _decodeDirection(int raw, int ewDir)\n"
         "  => ewDir == 1 ? raw + 180.0 : raw.toDouble();\n\n"
         "// SpeedMult で低速/高速域切替\n"
         "double _decodeSpeedH(int raw, int mult)\n"
         "  => mult == 1 ? raw * 0.75 : raw * 0.25;\n\n"
         "// 垂直速度 int8 × 0.5 m/s\n"
         "double _decodeSpeedV(int raw) {\n"
         "  final s = raw >= 128 ? raw - 256 : raw;\n"
         "  return s * 0.5;\n"
         "}",
         Inches(5.5), Inches(1.42), Inches(4.15), Inches(3.05), size=10)

    txb(s, "追加された新フィールド",
        M, Inches(4.58), IW, Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["フィールド", "内容", "取得元"],
          [
              ["altitudeMeters",    "気圧高度 QFE（地上=0m）",  "AltitudeBaro Byte13-14"],
              ["altitudeGeoMeters", "GPS WGS84 高度",           "AltitudeGeo  Byte15-16  ← 新規"],
              ["speedH",           "水平速度 m/s",              "SpeedH Byte3 + SpeedMult"],
              ["speedV",           "垂直速度 m/s（正=上昇）",   "SpeedV Byte4 (int8)"],
              ["direction",        "進行方向 degrees 0-359",    "Direction Byte2 + EWDir"],
              ["statusIsAirborne", "飛行中フラグ",              "Status == AIRBORNE(2)"],
          ],
          M, Inches(4.88), IW,
          col_widths=[Inches(2.0), Inches(3.0), Inches(4.1)], font_size=10)

    # ── スライド 18: Pi 起動エラーの修正 ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "Pi 起動エラーの修正",
           "hciconfig hci0 up — アダプターが既に UP のときにエラーになる問題")
    page_num(s, 18)

    txb(s, "発生したエラー",
        M, Inches(1.12), IW, Inches(0.28),
        size=11, bold=True, color=RGBColor(0xC6, 0x28, 0x28))
    code(s,
         "$ python3 -u /home/pi/remote_id.py --gps\n"
         "GPS: /dev/ttyACM0 を開きました\n"
         "BME280: I2C 0x76 地上気圧=1003.5hPa\n"
         "Traceback (most recent call last):\n"
         "  File \"/home/pi/remote_id.py\", line 415, in configure_adapter\n"
         "    _run([\"sudo\", \"hciconfig\", \"hci0\", \"up\"])\n"
         "subprocess.CalledProcessError: Command '['sudo', 'hciconfig', 'hci0', 'up']'\n"
         "    returned non-zero exit status 1.",
         M, Inches(1.42), IW, Inches(2.2), size=10)

    txb(s, "原因",
        M, Inches(3.72), IW, Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "# hci0 は起動済み（UP RUNNING）だった\n"
         "hci0: Type: Primary  Bus: UART\n"
         "      BD Address: D8:3A:DD:E2:55:36\n"
         "      UP RUNNING   ← 既に起動済み\n"
         "# hciconfig up は既に UP のときにエラーコード 1 を返す\n"
         "# subprocess.run(check=True) がそれを例外として扱っていた",
         M, Inches(4.02), IW, Inches(1.9), size=10)

    txb(s, "修正",
        M, Inches(6.02), Inches(4.3), Inches(0.28),
        size=11, bold=True, color=C_TEAL)
    code(s,
         "# 修正前: check=True で失敗\n"
         "_run([\"sudo\", \"hciconfig\", HCI_DEV, \"up\"])\n\n"
         "# 修正後: check なし（既に UP でも続行）\n"
         "subprocess.run(\n"
         "    [\"sudo\", \"hciconfig\", HCI_DEV, \"up\"],\n"
         "    stdout=DEVNULL, stderr=DEVNULL,\n"
         ")",
         M, Inches(6.32), Inches(4.3), Inches(0.9), size=10)

    # ── スライド 19: アプリ内コンソールログ ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    header(s, "アプリ内コンソールログ",
           "BLE 受信・センサー取得・ZK 証明の動きをアプリ画面で確認できるようにした")
    page_num(s, 19)

    txb(s, "追加した仕組み",
        M, Inches(1.12), Inches(4.3), Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    code(s,
         "// _log() ヘルパー: タイムスタンプ付きで蓄積\n"
         "void _log(String msg) {\n"
         "  final ts = '12:34:56';\n"
         "  setState(() {\n"
         "    _status = msg;       // 既存の状態バーも更新\n"
         "    _logs.add('[$ts] $msg');\n"
         "    if (_logs.length > 200) _logs.removeAt(0);\n"
         "  });\n"
         "  // 追加後に自動スクロール\n"
         "  _logScrollController.jumpTo(maxExtent);\n"
         "}",
         M, Inches(1.42), Inches(4.3), Inches(2.88), size=10)

    code(s,
         "// BLE パケット受信時にログ追加\n"
         "[BLE BasicID] UAS_ID=UAV-ZKP-M1-DEMO00 rssi=-62dBm\n"
         "[BLE Location] seq=4 lat=34.968581\n"
         "  lon=135.899214 alt=0.3m\n"
         "  airborne=true rssi=-62dBm\n\n"
         "// 証明フロー\n"
         "① Witness を計算しています...\n"
         "② Groth16 証明を生成しています...\n"
         "✓ 証明・検証 完了（2.3s）",
         Inches(5.5), Inches(1.42), Inches(4.15), Inches(2.88), size=10)

    txb(s, "ログタブの機能",
        M, Inches(4.42), IW, Inches(0.28),
        size=11, bold=True, color=C_NAVY)
    table(s,
          ["機能", "詳細"],
          [
              ["ターミナル風 UI",   "黒背景 / 等幅フォント / タイムスタンプ付き"],
              ["色分け表示",        "BLE受信=青 / 完了=緑 / エラー=赤 / その他=白"],
              ["自動スクロール",    "新しいログが来ると最下部に自動スクロール"],
              ["クリア機能",        "右上のゴミ箱アイコンでログを一括削除"],
              ["最大件数",          "200件まで保持、超えた場合は古いものから削除"],
          ],
          M, Inches(4.72), IW,
          col_widths=[Inches(2.0), Inches(7.1)], font_size=11)

    # ── スライド 15: まとめ ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_white_bg(s)
    rect(s, 0, 0, W, Inches(0.07), C_NAVY)
    rect(s, 0, Inches(0.07), Inches(0.07), H - Inches(0.07), C_BLUE)

    txb(s, "まとめ",
        M, Inches(0.15), IW, Inches(0.55),
        size=26, bold=True, color=C_NAVY)
    rect(s, M, Inches(0.72), Inches(1.8), Pt(2.5), C_TEAL)

    bullets(s, [
        "■  独自 BLE フォーマット → ASTM F3411-22a (OpenDroneID) 準拠に変更",
        "■  Service UUID 0xFFFA の Service Data に 25 バイト固定長メッセージを格納",
        "■  BasicID（機体識別）と Location（位置・高度・時刻）を 1 秒ごと交互に送信",
        "■  Flutter アプリ: UUID 検索 → Location パース → タイムスタンプ再構築 → 自動入力",
        "■  取得座標を正規化して Circom ZK 回路の private input に渡す",
        "■  公開するのは within30 フラグと Poseidon commitment のみ — 生座標は非公開",
        "■  Time Series 回路 (N=20) で録画中の全データポイントを一括証明",
        "■  GPS（USB CDC）+ BME280（I²C）を Raspberry Pi に接続し実機センサーを統合",
        "■  QFE 方式でドローン・ユーザー双方の高度を地上ゼロ基準に統一",
        "■  opendroneid-core-c 準拠: Status 上位ニブル / AIRBORNE=2 / BasicID Byte1 修正",
        "■  アプリ内コンソールログ: BLE受信・ZK証明の動きをリアルタイムで確認可能",
    ], M, Inches(0.95), IW, Inches(5.1), size=11)

    rect(s, 0, H - Inches(0.75), W, Inches(0.75), C_NAVY)
    txb(s, "Pi 送信 / Flutter 受信の両側が完全 opendroneid-core-c 準拠 — BLE 受信動作確認済み",
        M, H - Inches(0.65), IW, Inches(0.52),
        size=11, bold=True, color=RGBColor(0xB0, 0xBE, 0xFF),
        align=PP_ALIGN.CENTER)

    # ── 保存 ────────────────────────────────────────────────────────────────
    out = "OPENDRONEID_IMPLEMENTATION.pptx"
    prs.save(out)
    print(f"✅ Saved: {out}  ({len(prs.slides)} slides, 4:3, 白ベース, 游ゴシック)")


if __name__ == "__main__":
    build()
