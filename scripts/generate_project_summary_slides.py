from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("project_summary_2026-03-27_4x3.pptx")


SLIDES = [
    {
        "title": "UAV安全距離証明アプリ",
        "subtitle": "実装・統合・検証の要約",
        "bullets": [
            "Circom / Groth16 / rapidsnark 基盤",
            "Flutter / BLE / GPS / 気圧センサー統合",
            "動画証拠・ML Kit・YOLO 拡張",
        ],
    },
    {
        "title": "目的",
        "bullets": [
            "30m以内・5秒以内の接近判定回路",
            "スマホ内での証明生成・検証フロー",
            "UAVログと第三者端末ログの統合",
            "研究ユースケース寄りUI",
        ],
    },
    {
        "title": "ZKP基盤実装",
        "bullets": [
            "緯度・経度・高度・時刻入力",
            "平面近似3次元距離判定回路",
            "Groth16 セットアップ",
            ".wcd / .zkey / verification key 生成",
            "witness長不一致の解消",
        ],
    },
    {
        "title": "アプリ統合",
        "bullets": [
            "Flutter から witness 計算",
            "rapidsnark による証明生成",
            "端末内での検証実行",
            "成功結果・公開信号表示",
            "Pixel 8 実機動作確認",
        ],
    },
    {
        "title": "データ取得拡張",
        "bullets": [
            "GPS による第三者位置取得",
            "気圧センサーによる高度推定",
            "BLE によるUAVテレメトリ受信",
            "ドローン側ログ自動入力",
            "スマホ完結フロー",
        ],
    },
    {
        "title": "証拠・検知機能",
        "bullets": [
            "BLE受信後の動画撮影",
            "撮影時刻とZKP入力の対応付け",
            "ML Kit による軽量動画解析",
            "YOLO ライブ検知画面",
            "UAV専用モデル差し替え前提",
        ],
    },
    {
        "title": "成果と残課題",
        "bullets": [
            "証明生成・検証の端末内成立",
            "UAV受信から証明までの構成完成",
            "実機E2E疎通の最終確認待ち",
            "Raspberry Pi BLE 実運用確認",
            "UAV専用検知モデルの高精度化",
        ],
    },
]


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 246, 240)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.35))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(36, 74, 64)
    accent.line.fill.background()


def set_text_style(paragraph, font_size, bold=False, color=RGBColor(34, 34, 34)):
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.name = "Hiragino Sans"
    paragraph.font.color.rgb = color


def add_title_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.6), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec["title"]
    p.alignment = PP_ALIGN.LEFT
    set_text_style(p, 24, bold=True, color=RGBColor(25, 25, 25))

    subtitle_box = slide.shapes.add_textbox(Inches(0.72), Inches(2.05), Inches(8.0), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec["subtitle"]
    set_text_style(p, 13, color=RGBColor(70, 70, 70))

    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(8.2), Inches(3.5))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(spec["bullets"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"・{bullet}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        set_text_style(p, 18)
        p.space_after = Pt(12)


def add_content_slide(prs, spec, page_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(8.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec["title"]
    set_text_style(p, 22, bold=True, color=RGBColor(28, 28, 28))

    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.38), Inches(2.2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(196, 125, 68)
    line.line.fill.background()

    body_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.8), Inches(8.0), Inches(4.9))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(spec["bullets"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"・{bullet}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        set_text_style(p, 20)
        p.space_after = Pt(14)

    page_box = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(0.5), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_number)
    p.alignment = PP_ALIGN.RIGHT
    set_text_style(p, 10, color=RGBColor(100, 100, 100))


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, SLIDES[0])
    for idx, slide_spec in enumerate(SLIDES[1:], start=2):
        add_content_slide(prs, slide_spec, idx)

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
